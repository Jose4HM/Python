# Importar la biblioteca necesaria
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_slide(prs, title, content):
    # Función para agregar una diapositiva con título y contenido
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    content_box = slide.placeholders[1]
    content_frame = content_box.text_frame

    for paragraph in content.split('\n'):
        p = content_frame.add_paragraph()
        p.text = paragraph
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.LEFT

def create_presentation():
    # Función principal para crear la presentación
    prs = Presentation()

    # Introducción
    add_slide(prs, "Evolución y Desafíos de la Distribución de Contenido Audiovisual",
              "Tradicionalmente, la distribución de contenido a través de la televisión por cable ha dominado. "
              "Ingeniería de telecomunicaciones enfrenta limitaciones en ancho de banda y mantenimiento.")

    # Alternativas Tradicionales
    add_slide(prs, "Alternativas Tradicionales",
              "1. **Televisión Satelital:**\n"
              "   - Retardos significativos y afectación por condiciones climáticas.\n"
              "   - Amplia cobertura, pero inversiones sustanciales y mantenimiento meticuloso.\n"
              "2. **Difusión Terrestre:**\n"
              "   - Limitaciones en capacidad y calidad, especialmente en áreas remotas.\n"
              "   - Mejora con estándares digitales, pero desafíos persisten.")

    # Ventajas de la Televisión Streaming
    add_slide(prs, "Ventajas de la Televisión Streaming",
              "- Acceso a la carta y exploración de contenido bajo demanda.\n"
              "- Flexibilidad horaria, ubicación y comodidad para el usuario.")

    # Tecnología detrás de la Experiencia del Usuario
    add_slide(prs, "Tecnología detrás de la Experiencia del Usuario",
              "1. **Interfaz de Usuario y Navegación:**\n"
              "   - Eficiencia en navegación mediante transmisión rápida de datos.\n"
              "2. **Servidores y Almacenamiento:**\n"
              "   - Almacenamiento en la nube para carga rápida y accesibilidad.\n"
              "3. **Transmisión de Datos:**\n"
              "   - Velocidad, estabilidad y adaptación a las condiciones de la red.")

    # ... Agregar más diapositivas según sea necesario ...

    # Conclusiones
    add_slide(prs, "Conclusiones",
              "- Streaming rompe barreras de programación tradicionales.\n"
              "- Ingeniería de telecomunicaciones es esencial en toda la cadena de valor.\n"
              "- Recopilación y análisis de datos permiten ajustes continuos.\n"
              "- Seguridad y eficiencia son fundamentales para el éxito sostenido.")

    prs.save("Presentacion_Distribucion_Contenido.pptx")

if __name__ == "__main__":
    create_presentation()
